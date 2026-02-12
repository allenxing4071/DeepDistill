"""
文档处理器：PDF / Word / PPT / Excel / HTML → 文本
保留文档结构（标题层级、段落、表格）。

依赖：pip install deepdistill[doc]
"""

from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger("deepdistill.document")


def extract_text_from_document(file_path: Path) -> str:
    """根据文件类型分发到对应的提取器"""
    suffix = file_path.suffix.lower()

    if suffix == ".pdf":
        return _extract_pdf(file_path)
    elif suffix in (".docx", ".doc"):
        return _extract_docx(file_path)
    elif suffix in (".pptx", ".ppt"):
        return _extract_pptx(file_path)
    elif suffix in (".xlsx", ".xls"):
        return _extract_xlsx(file_path)
    elif suffix in (".txt", ".md"):
        return _extract_text_file(file_path)
    else:
        raise ValueError(f"不支持的文档格式: {suffix}")


def extract_text_from_html(file_path: Path) -> str:
    """从 HTML 文件提取正文"""
    return _extract_html(file_path)


def _extract_pdf(file_path: Path) -> str:
    """PDF → 文本"""
    from PyPDF2 import PdfReader

    logger.info(f"提取 PDF: {file_path.name}")
    reader = PdfReader(str(file_path))

    texts = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text:
            texts.append(f"--- 第 {i + 1} 页 ---\n{text.strip()}")

    full_text = "\n\n".join(texts)
    logger.info(f"PDF 提取完成: {len(reader.pages)} 页, {len(full_text)} 字符")
    return full_text


def _extract_docx(file_path: Path) -> str:
    """Word (.docx) → 文本"""
    from docx import Document

    logger.info(f"提取 Word: {file_path.name}")
    doc = Document(str(file_path))

    texts = []
    for para in doc.paragraphs:
        if para.text.strip():
            # 保留标题层级
            if para.style and para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading ", "")
                prefix = "#" * int(level) if level.isdigit() else "#"
                texts.append(f"{prefix} {para.text.strip()}")
            else:
                texts.append(para.text.strip())

    full_text = "\n\n".join(texts)
    logger.info(f"Word 提取完成: {len(texts)} 段, {len(full_text)} 字符")
    return full_text


def _extract_pptx(file_path: Path) -> str:
    """PPT (.pptx) → 文本"""
    from pptx import Presentation

    logger.info(f"提取 PPT: {file_path.name}")
    prs = Presentation(str(file_path))

    texts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())
        if slide_texts:
            texts.append(f"--- 幻灯片 {i + 1} ---\n" + "\n".join(slide_texts))

    # 提取备注
    for i, slide in enumerate(prs.slides):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                texts.append(f"--- 幻灯片 {i + 1} 备注 ---\n{notes}")

    full_text = "\n\n".join(texts)
    logger.info(f"PPT 提取完成: {len(prs.slides)} 页, {len(full_text)} 字符")
    return full_text


def _extract_xlsx(file_path: Path) -> str:
    """Excel (.xlsx) → 文本"""
    from openpyxl import load_workbook

    logger.info(f"提取 Excel: {file_path.name}")
    wb = load_workbook(str(file_path), data_only=True)

    texts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            texts.append(f"--- 工作表: {sheet_name} ---\n" + "\n".join(rows))

    full_text = "\n\n".join(texts)
    logger.info(f"Excel 提取完成: {len(wb.sheetnames)} 表, {len(full_text)} 字符")
    return full_text


def _extract_text_file(file_path: Path) -> str:
    """纯文本 / Markdown → 文本"""
    logger.info(f"读取文本文件: {file_path.name}")
    return file_path.read_text(encoding="utf-8")


def _extract_html(file_path: Path) -> str:
    """HTML → 正文提取"""
    from bs4 import BeautifulSoup

    logger.info(f"提取 HTML: {file_path.name}")
    html = file_path.read_text(encoding="utf-8")
    soup = BeautifulSoup(html, "html.parser")

    # 移除脚本和样式
    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    text = soup.get_text(separator="\n", strip=True)
    logger.info(f"HTML 提取完成: {len(text)} 字符")
    return text
